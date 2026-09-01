"""Motor de presentación: un mismo spec -> PPTX (python-pptx) + preview PNG (Pillow).

Permite revisar visualmente cada lámina antes de entregar el archivo final.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont
import copy, os

# ─────────────────────────── design tokens ───────────────────────────
W, H = 20.0, 11.25                      # pulgadas (16:9)
MX = 1.05                               # margen lateral
CW = W - 2 * MX                         # ancho de contenido = 17.9

ROJO   = "D10120"
AMAR   = "F7B006"
AMAR_D = "B07800"
AZUL   = "005F91"
AZUL_D = "003F62"
INK    = "13161B"
TXT    = "394049"
GRAY   = "6E6E6E"
MUT    = "9AA1A9"
LIGHT  = "F5F7F9"
LIGHT2 = "EFF3F6"
LINE   = "E2E6EA"
WHITE  = "FFFFFF"

FONT = "Open Sans"                      # tipografía del PPTX

# geometría de lámina estándar
EYE_Y, TIT_Y, RULE_Y = 0.66, 1.00, 2.06
CT, CB = 2.46, 10.02                    # content top / bottom
FOOT_Y = 10.42

FONTS_DIR = "/home/user/work/fonts/ssp/font_source_sans_pro/files"
_PREVIEW_FONT = {
    (False, False): f"{FONTS_DIR}/SourceSansPro-Regular.ttf",
    (True, False):  f"{FONTS_DIR}/SourceSansPro-Bold.ttf",
    (False, True):  f"{FONTS_DIR}/SourceSansPro-It.ttf",
    (True, True):   f"{FONTS_DIR}/SourceSansPro-BoldIt.ttf",
}
_SEMI = f"{FONTS_DIR}/SourceSansPro-Semibold.ttf"
_fcache = {}
DPI = 96
WIDTH_SAFETY = 1.06        # Open Sans es ~6% más ancha que Source Sans Pro


def _pf(size_pt, bold=False, italic=False, semi=False):
    key = (round(size_pt, 1), bold, italic, semi)
    if key not in _fcache:
        path = _SEMI if semi else _PREVIEW_FONT[(bold, italic)]
        _fcache[key] = ImageFont.truetype(path, int(round(size_pt * DPI / 72)))
    return _fcache[key]


def rgb(h):
    return RGBColor.from_string(h)


def expand(paras):
    """Convierte saltos de línea explícitos en párrafos independientes."""
    out = []
    for p in paras:
        parts = p["text"].split("\n")
        for i, t in enumerate(parts):
            q = dict(p)
            q["text"] = t
            if i < len(parts) - 1:
                q["space_after"] = 0
            if i > 0:
                q["space_before"] = 0
            out.append(q)
    return out


def para(text, size=14.5, bold=False, color=TXT, align="l", space_before=0, space_after=6,
         line=1.16, bullet=False, italic=False, semi=False, spc=0, indent=0.30):
    return dict(text=text, size=size, bold=bold, color=color, align=align,
                space_before=space_before, space_after=space_after, line=line,
                bullet=bullet, italic=italic, semi=semi, spc=spc, indent=indent)


# ─────────────────────────── slide ───────────────────────────
class Slide:
    def __init__(self, deck, number=None):
        self.deck = deck
        self.ops = []
        self.number = number
        self.notes = None

    # primitivas -----------------------------------------------------
    def rect(self, x, y, w, h, fill=None, line=None, lw=0.75, radius=0.0):
        self.ops.append(("rect", dict(x=x, y=y, w=w, h=h, fill=fill, line=line, lw=lw, radius=radius)))

    def hline(self, x, y, w, color=LINE, lw=0.75):
        self.ops.append(("line", dict(x1=x, y1=y, x2=x + w, y2=y, color=color, lw=lw)))

    def vline(self, x, y, h, color=LINE, lw=0.75):
        self.ops.append(("line", dict(x1=x, y1=y, x2=x, y2=y + h, color=color, lw=lw)))

    def arrow(self, x1, y1, x2, y2, color=MUT, lw=1.5):
        self.ops.append(("arrow", dict(x1=x1, y1=y1, x2=x2, y2=y2, color=color, lw=lw)))

    def text(self, x, y, w, h, paras, valign="t"):
        if isinstance(paras, dict):
            paras = [paras]
        self.ops.append(("text", dict(x=x, y=y, w=w, h=h, paras=paras, valign=valign)))

    def image(self, path, x, y, w, h, align="c", valign="c"):
        """Encaja la imagen dentro de la caja (x,y,w,h) preservando proporción."""
        iw, ih = Image.open(path).size
        ar = iw / ih
        bw, bh = w, h
        if bw / bh > ar:
            nh, nw = bh, bh * ar
        else:
            nw, nh = bw, bw / ar
        nx = x + {"l": 0, "c": (bw - nw) / 2, "r": bw - nw}[align]
        ny = y + {"t": 0, "c": (bh - nh) / 2, "b": bh - nh}[valign]
        self.ops.append(("img", dict(path=path, x=nx, y=ny, w=nw, h=nh)))
        return (nx, ny, nw, nh)

    # componentes ----------------------------------------------------
    def header(self, eyebrow_num, eyebrow, title, subtitle=None):
        if eyebrow or eyebrow_num:
            self.text(MX, EYE_Y, CW, 0.34, [
                para(f"{eyebrow_num}   {eyebrow}".strip(), size=13.5, bold=True, color=ROJO,
                     space_after=0, spc=2.4)], valign="t")
        self.text(MX, TIT_Y, CW, 0.92, [para(title, size=33, bold=True, color=INK,
                                             space_after=0, line=1.05)], valign="t")
        for i, c in enumerate((ROJO, AMAR, AZUL)):
            self.rect(MX + i * 0.92, RULE_Y, 0.86, 0.085, fill=c)
        if subtitle:
            self.text(MX, RULE_Y + 0.24, CW * 0.86, 0.5,
                      [para(subtitle, size=16.5, color=GRAY, space_after=0, line=1.2)])

    def footer(self, text="Pronóstico Jerárquico de Demanda Intermitente en el Retail de Moda Femenina"):
        self.hline(MX, FOOT_Y - 0.16, CW, color=LINE, lw=0.75)
        self.text(MX, FOOT_Y, CW * 0.8, 0.3,
                  [para(text, size=10.5, color=MUT, space_after=0)])
        if self.number:
            self.text(MX + CW - 1.4, FOOT_Y, 1.4, 0.3,
                      [para(f"{self.number:02d}", size=11.5, bold=True, color=AZUL,
                            align="r", space_after=0)])

    def card(self, x, y, w, h, title=None, body=None, accent=None, fill=LIGHT,
             tsize=17, bsize=14.2, pad=0.42, border=None, title_color=INK,
             accent_style="left"):
        self.rect(x, y, w, h, fill=fill, line=border, radius=0.055)
        side = accent and accent_style == "left"
        if side:
            self.rect(x, y, 0.075, h, fill=accent, radius=0.0)
        elif accent and accent_style == "top":
            self.rect(x, y, w, 0.075, fill=accent, radius=0.0)
        tx = x + pad + (0.08 if side else 0)
        tw = w - 2 * pad - (0.08 if side else 0)
        paras = []
        if title:
            paras.append(para(title, size=tsize, bold=True, color=title_color,
                              space_after=8 if body else 0, line=1.12))
        if body:
            paras += body
        self.text(tx, y + pad - 0.06, tw, h - 2 * pad + 0.12, paras)

    def stat(self, x, y, w, h, value, label, color=AZUL, fill=LIGHT, vsize=46, lsize=13.5,
             note=None, align="l"):
        self.rect(x, y, w, h, fill=fill, radius=0.055)
        pad = 0.36
        paras = [para(value, size=vsize, bold=True, color=color, space_after=2, line=1.0, align=align),
                 para(label, size=lsize, color=TXT, space_after=0, line=1.2, align=align)]
        if note:
            paras.append(para(note, size=11.5, color=MUT, space_before=4, space_after=0, line=1.15, align=align))
        self.text(x + pad, y, w - 2 * pad, h, paras, valign="m")

    def pill(self, x, y, w, h, text, fill=AZUL, color=WHITE, size=14, bold=True, radius=0.5):
        self.rect(x, y, w, h, fill=fill, radius=radius)
        self.text(x + 0.16, y, w - 0.32, h, [para(text, size=size, bold=bold, color=color,
                                                  align="c", space_after=0, line=1.1)], valign="m")

    def takeaway(self, x, y, w, text, color=ROJO, size=16.5):
        """Cierre de lámina sin caja: regla corta de color + frase. Reemplaza al banner."""
        self.rect(x, y, 0.62, 0.075, fill=color)
        h = max(0.5, min(1.3, FOOT_Y - 0.22 - (y + 0.26)))
        self.text(x, y + 0.26, w, h,
                  [para(text, size=size, bold=True, color=INK, space_after=0, line=1.28)])

    def caption(self, x, y, w, text, align="l", size=11.5):
        self.text(x, y, w, 0.42, [para(text, size=size, color=MUT, align=align, space_after=0, line=1.2)])

    def banner(self, x, y, w, h, text, fill=AZUL, color=WHITE, size=16, bold=True, align="c"):
        self.rect(x, y, w, h, fill=fill, radius=0.055)
        self.text(x + 0.5, y, w - 1.0, h, [para(text, size=size, bold=bold, color=color,
                                                align=align, space_after=0, line=1.18)], valign="m")


# ─────────────────────────── deck ───────────────────────────
class Deck:
    def __init__(self):
        self.slides = []

    def new(self, number=None):
        s = Slide(self, number)
        self.slides.append(s)
        return s

    # ---------------- PPTX ----------------
    def save_pptx(self, path):
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(W), Inches(H)
        blank = prs.slide_layouts[6]
        for sl in self.slides:
            slide = prs.slides.add_slide(blank)
            for kind, o in sl.ops:
                getattr(self, f"_pptx_{kind}")(slide, o)
            if sl.notes:
                tf = slide.notes_slide.notes_text_frame
                tf.text = sl.notes
        cp = prs.core_properties
        cp.title = "Pronóstico Jerárquico de Demanda Intermitente en el Retail de Moda Femenina"
        cp.author = "Joaquín Ignacio Mondaca Parada"
        cp.subject = "Memoria para optar al título de Ingeniero Comercial — UTFSM"
        prs.save(path)

    def _pptx_rect(self, slide, o):
        shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if o["radius"] > 0 else MSO_SHAPE.RECTANGLE
        s = slide.shapes.add_shape(shp_type, Inches(o["x"]), Inches(o["y"]),
                                   Inches(o["w"]), Inches(o["h"]))
        if o["radius"] > 0:
            try:
                s.adjustments[0] = min(0.5, o["radius"] / min(o["w"], o["h"]))
            except Exception:
                pass
        if o["fill"]:
            s.fill.solid(); s.fill.fore_color.rgb = rgb(o["fill"])
        else:
            s.fill.background()
        if o["line"]:
            s.line.color.rgb = rgb(o["line"]); s.line.width = Pt(o["lw"])
        else:
            s.line.fill.background()
        s.shadow.inherit = False
        s.text_frame.word_wrap = False
        return s

    def _pptx_line(self, slide, o):
        from pptx.enum.shapes import MSO_CONNECTOR
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(o["x1"]), Inches(o["y1"]),
                                       Inches(o["x2"]), Inches(o["y2"]))
        c.line.color.rgb = rgb(o["color"]); c.line.width = Pt(o["lw"])
        return c

    def _pptx_arrow(self, slide, o):
        c = self._pptx_line(slide, o)
        ln = c.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
        return c

    def _pptx_img(self, slide, o):
        slide.shapes.add_picture(o["path"], Inches(o["x"]), Inches(o["y"]),
                                 Inches(o["w"]), Inches(o["h"]))

    def _pptx_text(self, slide, o):
        box = slide.shapes.add_textbox(Inches(o["x"]), Inches(o["y"]), Inches(o["w"]), Inches(o["h"]))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                              "b": MSO_ANCHOR.BOTTOM}[o["valign"]]
        for i, p in enumerate(expand(o["paras"])):
            pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            pp.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                            "r": PP_ALIGN.RIGHT, "j": PP_ALIGN.JUSTIFY}[p["align"]]
            pp.line_spacing = p["line"]
            pp.space_before = Pt(p["space_before"])
            pp.space_after = Pt(p["space_after"])
            if p["bullet"]:
                pPr = pp._p.get_or_add_pPr()
                pPr.set("marL", str(int(p["indent"] * 914400)))
                pPr.set("indent", str(int(-p["indent"] * 914400)))
            txt = ("•\t" if p["bullet"] else "") + p["text"]
            r = pp.add_run(); r.text = txt
            f = r.font
            f.name = FONT + (" SemiBold" if p["semi"] else "")
            f.size = Pt(p["size"]); f.bold = p["bold"]; f.italic = p["italic"]
            f.color.rgb = rgb(p["color"])
            if p["spc"]:
                r.font._rPr.set("spc", str(int(p["spc"] * 100)))
        return box

    # ---------------- PREVIEW ----------------
    def render_previews(self, outdir, only=None):
        os.makedirs(outdir, exist_ok=True)
        paths = []
        for i, sl in enumerate(self.slides, 1):
            if only and i not in only:
                continue
            img = Image.new("RGB", (int(W * DPI), int(H * DPI)), "white")
            d = ImageDraw.Draw(img)
            for kind, o in sl.ops:
                getattr(self, f"_pv_{kind}")(img, d, o)
            p = f"{outdir}/slide{i:02d}.png"
            img.save(p); paths.append(p)
        return paths

    @staticmethod
    def _px(v):
        return v * DPI

    def _pv_rect(self, img, d, o):
        x, y, w, h = [self._px(o[k]) for k in ("x", "y", "w", "h")]
        r = self._px(o["radius"])
        box = [x, y, x + w, y + h]
        fill = "#" + o["fill"] if o["fill"] else None
        outline = "#" + o["line"] if o["line"] else None
        lw = max(1, int(o["lw"] * DPI / 72))
        if r > 0:
            d.rounded_rectangle(box, radius=min(r, min(w, h) / 2), fill=fill, outline=outline, width=lw)
        else:
            d.rectangle(box, fill=fill, outline=outline, width=lw)

    def _pv_line(self, img, d, o):
        d.line([self._px(o["x1"]), self._px(o["y1"]), self._px(o["x2"]), self._px(o["y2"])],
               fill="#" + o["color"], width=max(1, int(o["lw"] * DPI / 72)))

    def _pv_arrow(self, img, d, o):
        self._pv_line(img, d, o)
        import math
        x1, y1, x2, y2 = [self._px(o[k]) for k in ("x1", "y1", "x2", "y2")]
        ang = math.atan2(y2 - y1, x2 - x1)
        L = 9
        pts = [(x2, y2),
               (x2 - L * math.cos(ang - 0.5), y2 - L * math.sin(ang - 0.5)),
               (x2 - L * math.cos(ang + 0.5), y2 - L * math.sin(ang + 0.5))]
        d.polygon(pts, fill="#" + o["color"])

    def _pv_img(self, img, d, o):
        im = Image.open(o["path"])
        if im.mode == "RGBA":
            bg = Image.new("RGB", im.size, "white"); bg.paste(im, mask=im.split()[3]); im = bg
        im = im.convert("RGB").resize((max(1, int(self._px(o["w"]))), max(1, int(self._px(o["h"])))),
                                      Image.LANCZOS)
        img.paste(im, (int(self._px(o["x"])), int(self._px(o["y"]))))

    def _pv_text(self, img, d, o):
        lines = self.layout_text(o)
        total = sum(l["h"] for l in lines)
        boxh = self._px(o["h"])
        y0 = self._px(o["y"])
        if o["valign"] == "m":
            y0 += (boxh - total) / 2
        elif o["valign"] == "b":
            y0 += boxh - total
        x0, boxw = self._px(o["x"]), self._px(o["w"])
        for l in lines:
            f = l["font"]
            tw = d.textlength(l["text"], font=f)
            if l["align"] == "c":
                tx = x0 + (boxw - tw) / 2
            elif l["align"] == "r":
                tx = x0 + boxw - tw
            else:
                tx = x0 + l["dx"]
            d.text((tx, y0 + l["pad_top"]), l["text"], font=f, fill="#" + l["color"])
            y0 += l["h"]

    def layout_text(self, o):
        """Devuelve líneas con métricas (usado por preview y por el control de calidad)."""
        out = []
        boxw = self._px(o["w"])
        for p in expand(o["paras"]):
            f = _pf(p["size"], p["bold"], p["italic"], p["semi"])
            lh = p["size"] * 1.32 * p["line"] * DPI / 72
            indent = self._px(p["indent"]) if p["bullet"] else 0
            words = p["text"].split(" ")
            cur, curdx, first = "", 0, True
            plines = []

            def avail(first_line):
                return boxw - (0 if first_line else indent)

            dummy = ImageDraw.Draw(Image.new("RGB", (1, 1)))
            prefix = "•     " if p["bullet"] else ""
            for wd in words:
                trial = (cur + " " + wd).strip()
                test = (prefix + trial) if (first and p["bullet"]) else trial
                if dummy.textlength(test, font=f) * WIDTH_SAFETY <= avail(first) or not cur:
                    cur = trial
                else:
                    plines.append((prefix + cur if first and p["bullet"] else cur, 0 if first else indent))
                    first = False
                    cur = wd
            plines.append((prefix + cur if first and p["bullet"] else cur, 0 if first else indent))
            for i, (t, dx) in enumerate(plines):
                out.append(dict(text=t, font=f, color=p["color"], align=p["align"], dx=dx, h=lh,
                                pad_top=(p["space_before"] * DPI / 72 if i == 0 else 0)))
            if out:
                out[-1]["h"] += p["space_after"] * DPI / 72
                if p["space_before"] and len(plines):
                    out[-len(plines)]["h"] += p["space_before"] * DPI / 72
        return out

    def save_pdf(self, path, outdir="/home/user/work/preview"):
        """PDF de respaldo generado desde los previews (misma maqueta)."""
        from PIL import Image
        imgs = [Image.open(p).convert("RGB") for p in sorted(__import__("glob").glob(f"{outdir}/slide*.png"))]
        if imgs:
            imgs[0].save(path, save_all=True, append_images=imgs[1:], resolution=DPI)
        return path

    # ---------------- QA ----------------
    def qa(self):
        from fontTools.ttLib import TTFont
        cov = set()
        for f in (_PREVIEW_FONT[(False, False)], _PREVIEW_FONT[(True, False)]):
            cov |= set(TTFont(f).getBestCmap().keys())
        issues = []
        for i, sl in enumerate(self.slides, 1):
            for kind, o in sl.ops:
                if kind == "text":
                    for p in o["paras"]:
                        missing = {c for c in p["text"] if ord(c) not in cov and c not in " \n\t"}
                        if missing:
                            issues.append(f"S{i:02d} GLIFO AUSENTE {sorted(missing)} :: {p['text'][:50]!r}")
        for i, sl in enumerate(self.slides, 1):
            for kind, o in sl.ops:
                if kind == "text":
                    lines = self.layout_text(o)
                    need = sum(l["h"] for l in lines) / DPI
                    if need > o["h"] + 0.02:
                        issues.append(f"S{i:02d} TEXTO DESBORDA {need:.2f}in > caja {o['h']:.2f}in :: "
                                      f"{o['paras'][0]['text'][:55]!r}")
                    if o["x"] < 0.3 or o["x"] + o["w"] > W - 0.3:
                        issues.append(f"S{i:02d} TEXTO fuera de márgenes x={o['x']:.2f} w={o['w']:.2f}")
                if kind in ("rect", "img", "text"):
                    if o["y"] + o["h"] > FOOT_Y - 0.20 and not (kind == "text" and o["h"] < 0.45):
                        issues.append(f"S{i:02d} {kind.upper()} INVADE EL PIE  y+h={o['y']+o['h']:.2f} > 10.22")
                for k in ("rect", "img"):
                    if kind == k:
                        if o["y"] + o["h"] > H - 0.15 or o["y"] < 0.1:
                            issues.append(f"S{i:02d} {kind.upper()} fuera de lámina y={o['y']:.2f} h={o['h']:.2f}")
        return issues
